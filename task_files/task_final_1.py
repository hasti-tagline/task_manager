from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
import re
from pptx.oxml.ns import qn


prs = Presentation()

def setup_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.1),
        prs.slide_width - Inches(0.5),
        Inches(1)
    )

    p = title_box.text_frame.paragraphs[0]
    p.text = "THEMES BY ANALYSIS GROUPS"
    p.alignment = PP_ALIGN.RIGHT

    for r in p.runs:
        r.font.size = Pt(32)
        r.font.bold = True
        r.font.color.rgb = RGBColor(230, 170, 0)

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(245, 235, 210)

    return slide

# -----------------------
# LEGEND
# -----------------------
def add_legend(slide):

    legend_data = [
        ("Much less than total\n(-11 or less)", RGBColor(230, 85, 13)),
        ("Less than total\n(-4 to -10)", RGBColor(253, 141, 60)),
        ("In line with total\n(-3 to +3)", RGBColor(220, 220, 220)),
        ("More than total\n(+4 to +10)", RGBColor(66, 133, 244)),
        ("Much more than total\n(+11 or more)", RGBColor(0, 51, 153))
    ]

    legend_left = Inches(3.2)
    available_width = prs.slide_width - legend_left
    gap = Inches(0.05)
    box_width = (available_width - gap * 4) / 5
    top = Inches(0.9)

    for i, (text, color) in enumerate(legend_data):
        left = legend_left + i * (box_width + gap)

        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, box_width, Inches(0.5)
        )

        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER

        for r in p.runs:
            r.font.size = Pt(10)
            r.font.bold = True
            r.font.color.rgb = RGBColor(255,255,255) if i != 2 else RGBColor(0,0,0)

def set_border(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for line in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = OxmlElement(line)

        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '000000')  # black border
        solidFill.append(srgbClr)

        ln.append(solidFill)
        tcPr.append(ln)

def get_blue(percent):
    base = (7, 27, 69)  # your desired color

    ratio = (percent - 60) / 35
    factor = 0.5 + ratio * 0.5

    r = int(base[0] * factor + 255 * (1 - factor))
    g = int(base[1] * factor + 255 * (1 - factor))
    b = int(base[2] * factor + 255 * (1 - factor))

    return RGBColor(r, g, b)
# -----------------------
# DATA
# -----------------------
group_definitions = {
    "FREQUENCY": ["0 times","1-2 times","3-5 times","6-7 times","8+ times"],
    "INCOME": ["High Income","Mid Income","Low Income"],
    "GENDER": ["Male","Female"],
    "IMPORTANCE": ["Very","Somewhat","Not"],
    "LOCATION": ["Urban","Rural"],
    "INCOME": ["High Income","Mid Income","Low Income"]
}

themes = [
    "VARIETY OF ACTIVITIES AND ENTERTAINMENT 92%",
    "BEAUTIFUL BEACHES AND NATURE 89%",
    "SPORTS EVENTS 76%",
    "CULTURAL DIVERSITY 71%",
    "PLEASANT WEATHER 68%",
    "FRIENDLY COMMUNITY 65%"
]

data = [
    [3,-1,5,0,2,"N/A",7,1,6,-2,4,8,0,1,-2],
    [-2,4,1,6,"N/A",-3,5,2,0,8,-1,3,7,-11,15],
    [0,6,-2,3,5,1,"N/A",-4,7,2,9,0,4,20,-5],
    [-5,2,4,1,0,"N/A",-2,6,3,8,7,5,1,0,"N/A"],
    [7,"N/A",-3,5,2,8,-1,4,6,0,3,9,2,-6,-1],
    [1,-4,6,2,7,"N/A",3,5,0,-2,8,4,6,0,12]
]

columns = [
    "0 times","1-2 times","3-5 times","6-7 times","8+ times",
    "High Income","Mid Income","Low Income",
    "Male","Female",
    "Very","Somewhat","Not",
    "Urban","Rural",

]

MAX_ROWS = 12
MAX_COLS = 10


CELL_HEIGHT = Inches(0.35)
CELL_WIDTH = Inches(1.0)


MAX_CELL_HEIGHT = Inches(0.4)
MAX_THEME_WIDTH = Inches(3)
TABLE_TOP = Inches(1.5)
TABLE_HEIGHT = Inches(4.5)
TABLE_WIDTH = prs.slide_width

rows_data = len(themes)
cols_data = len(columns)

needs_split = rows_data > MAX_ROWS or cols_data > MAX_COLS

rows = rows_data + 2   # +2 for header rows
cols = cols_data + 1 

# -----------------------
# TABLE
# -----------------------


if not needs_split:
    # -----------------------
    # SLIDE
    # -----------------------
    slide = setup_slide(prs)
    add_legend(slide)

    table = slide.shapes.add_table(
    rows, cols,
    Inches(0), TABLE_TOP,
    table_width, table_height
    ).table
    table.allow_autofit = False
    table_height = CELL_HEIGHT * rows
    table_width = CELL_WIDTH * cols 

    for r in table.rows:
        r.height = CELL_HEIGHT

    for c in table.columns:
        c.width = CELL_WIDTH


    # Column widths
    theme_width = int(TABLE_WIDTH * 0.3)
    other_width = int((TABLE_WIDTH - theme_width) / (cols - 1))

    table.columns[0].width = theme_width

    for i in range(1, cols):
        table.columns[i].width = other_width

    # -----------------------
    # ✅ FIXED HEADERS (FINAL)
    # -----------------------

    # Map column → group
    col_to_group = {}
    for group_name, group_cols in group_definitions.items():
        for col in group_cols:
            col_to_group[col] = group_name

    col_index = 1
    start = 1
    current_group = col_to_group.get(columns[0], "")

    for i in range(len(columns)):
        group = col_to_group.get(columns[i], "")

        if group != current_group:
            end = col_index - 1

            if start <= end and start < cols:
                table.cell(0, start).text = current_group

                if start < end:
                    try:
                        table.cell(0, start).merge(table.cell(0, end))
                    except:
                        pass

            start = col_index
            current_group = group

        col_index += 1

    # LAST GROUP
    end = col_index - 1

    if start <= end and start < cols:
        table.cell(0, start).text = current_group

        if start < end:
            try:
                table.cell(0, start).merge(table.cell(0, end))
            except:
                pass

    # -----------------------
    # THEMES HEADER
    # -----------------------
    themes_cell = table.cell(0,0)
    themes_cell.merge(table.cell(1,0))
    themes_cell.text = "THEMES"

    # -----------------------
    # COLUMN LABELS
    # -----------------------
    for j, col in enumerate(columns):
        if j+1 >= cols:
            break
        table.cell(1, j+1).text = col

    # -----------------------
    # HEADER STYLE
    # -----------------------
    HEADER_YELLOW = RGBColor(255,200,0)

    for i in range(2):
        for j in range(cols):
            if i == 1 and j == 0:
                continue

            cell = table.cell(i,j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_YELLOW
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True
            cell.auto_size = None

            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(12)
                    r.font.color.rgb = RGBColor(0,0,0)

    # THEMES header bigger
    for r in themes_cell.text_frame.paragraphs[0].runs:
        r.font.size = Pt(16)

    # -----------------------
    # BORDER
    # -----------------------
    def set_border(cell):
        tcPr = cell._tc.get_or_add_tcPr()
        for b in ['lnL','lnR','lnT','lnB']:
            ln = OxmlElement(f'a:{b}')
            ln.set('w','20000')
            tcPr.append(ln)

    for row in table.rows:
        for cell in row.cells:
            set_border(cell)

    # -----------------------
    # BLUE SHADE FUNCTION
    # -----------------------
    def get_blue(percent):
        base = (7, 27, 69) 
        ratio = (percent-60)/35
        factor = 0.5 + ratio*0.5

        r = int(base[0]*factor + 255*(1-factor))
        g = int(base[1]*factor + 255*(1-factor))
        b = int(base[2]*factor + 255*(1-factor))

        return RGBColor(r,g,b)

    # -----------------------
    # FILL DATA
    # -----------------------
    LIGHT = RGBColor(235,225,200)

    for i, theme in enumerate(themes):
        row_idx = i + 2

        # THEMES COLUMN
        cell = table.cell(row_idx,0)
        cell.text = theme

        percent = int(re.search(r'(\d+)%', theme).group(1))
        cell.fill.solid()
        cell.fill.fore_color.rgb = get_blue(percent)

        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        cell.text_frame.margin_left = Inches(0.05)
        cell.text_frame.margin_top = Inches(0.05)
        cell.text_frame.margin_bottom = Inches(0.05)

        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = RGBColor(255,255,255)
                r.font.size = Pt(14)
                r.auto_size = None
                

        # DATA CELLS
        for j, val in enumerate(data[i]):
            if j+1 >= cols:
                break

            cell = table.cell(row_idx, j+1)
            cell.text = str(val)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.text_frame.word_wrap = True
            cell.word_wrap = True
            cell.auto_size = None

            if isinstance(val,str):
                cell.fill.fore_color.rgb = LIGHT
            elif val <= -11:
                cell.fill.fore_color.rgb = RGBColor(230,85,13)
            elif val <= -4:
                cell.fill.fore_color.rgb = RGBColor(253,141,60)
            elif val <= 3:
                cell.fill.fore_color.rgb = RGBColor(200,200,200)
            elif val <= 10:
                cell.fill.fore_color.rgb = RGBColor(66,133,244)
            else:
                cell.fill.fore_color.rgb = RGBColor(0,51,153)

            font_color = RGBColor(255,255,255)
            if isinstance(val, str) or (-3 <= val <= 3):
                font_color = RGBColor(0,0,0)

            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.color.rgb = font_color

else:

    # -----------------------
    # ✅ GROUP-BASED SPLIT (STRICT: NO GROUP BREAK)
    # -----------------------
    grouped_slides = []
    current_slide = []
    current_count = 0

    for group_name, group_cols in group_definitions.items():

        group_size = len(group_cols)

        # ❗ Case 1: group itself أكبر than MAX_COLS → split inside group
        if group_size > MAX_COLS:
            if current_slide:
                grouped_slides.append(current_slide)
                current_slide = []
                current_count = 0

            # split big group into chunks
            for i in range(0, group_size, MAX_COLS):
                chunk = group_cols[i:i+MAX_COLS]
                grouped_slides.append([(group_name, chunk)])

            continue

        # ❗ Case 2: adding group exceeds limit → new slide
        if current_count + group_size > MAX_COLS:
            grouped_slides.append(current_slide)
            current_slide = []
            current_count = 0

        # ✅ store group WITH columns
        current_slide.append((group_name, group_cols))
        current_count += group_size

    # last slide
    if current_slide:
        grouped_slides.append(current_slide)
    # -----------------------
    # ✅ LOOP THROUGH SLIDES
    # -----------------------

    for group_chunk in grouped_slides:

        col_subset = []
        for group_name, group_cols in group_chunk:
            col_subset.extend(group_cols)

        for row_start in range(0, len(themes), MAX_ROWS):

            theme_subset = themes[row_start:row_start + MAX_ROWS]

            # ✅ SAFE DATA MAPPING (NO INDEX ERROR)
            data_subset = [
                [row[columns.index(col)] for col in col_subset] 
                for row in data[row_start:row_start + MAX_ROWS]
            ]

            slide = setup_slide(prs)
            add_legend(slide)

            # -----------------------
            # TABLE
            # -----------------------
            rows = len(theme_subset) + 2
            cols = len(col_subset) + 1

            table_height = CELL_HEIGHT * rows
            table_width = CELL_WIDTH * cols 
            table = slide.shapes.add_table(
                rows, cols,
                Inches(0), TABLE_TOP,
                table_width, table_height
            ).table
            
            table.allow_autofit = False

    

            for r in table.rows:
                r.height = CELL_HEIGHT

            for c in table.columns:
                c.width = CELL_WIDTH

            # COLUMN WIDTH
            theme_width = int(TABLE_WIDTH * 0.3)
            other_width = int((TABLE_WIDTH - theme_width) / (cols - 1))

            table.columns[0].width = theme_width

            for i in range(1, cols):
                table.columns[i].width = other_width

            # -----------------------
            # HEADERS (SAFE)
            # -----------------------
            col_index = 1

            for group_name, group_cols in group_definitions.items():

                if col_index >= cols:
                    break

                matched_cols = [c for c in col_subset if c in group_cols]

                if not matched_cols:
                    continue

                start = col_index
                end = col_index + len(matched_cols) - 1

                # clamp
                end = min(end, cols - 1)

                if start >= cols:
                    break

                table.cell(0, start).text = group_name

                if start < end:
                    try:
                        table.cell(0, start).merge(table.cell(0, end))
                    except:
                        pass

                col_index = end + 1

            # THEMES HEADER
            themes_cell = table.cell(0,0)
            themes_cell.merge(table.cell(1,0))
            themes_cell.text = "THEMES"

            for j, col in enumerate(col_subset):
                table.cell(1, j+1).text = col

            # -----------------------
            # HEADER STYLE
            # -----------------------
            HEADER_YELLOW = RGBColor(255,200,0)

            for i in range(2):
                for j in range(cols):
                    if i==1 and j==0:
                        continue

                    cell = table.cell(i,j)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = HEADER_YELLOW
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    cell.word_wrap = True
                    cell.auto_size = None

                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        for r in p.runs:
                            r.font.bold = True
                            r.font.size = Pt(14)
                            r.font.color.rgb = RGBColor(0,0,0)


            # BORDER
            for row in table.rows:
                for cell in row.cells:
                    set_border(cell)

            # -----------------------
            # FILL DATA
            # -----------------------
            LIGHT = RGBColor(235,225,200)

            for i, theme in enumerate(theme_subset):
                row_idx = i + 2

                # THEMES COLUMN
                cell = table.cell(row_idx,0)
                cell.text = theme

                percent = int(re.search(r'(\d+)%', theme).group(1))
                cell.fill.solid()
                cell.fill.fore_color.rgb = get_blue(percent)

                cell.text_frame.margin_left = Inches(0.05)

                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    for r in p.runs:
                        r.font.color.rgb = RGBColor(255,255,255)
                        r.font.size = Pt(12)

                # DATA CELLS
                for j, val in enumerate(data_subset[i]):
                    cell = table.cell(row_idx, j+1)
                    cell.text = str(val)

                    cell.fill.solid()
                    cell.word_wrap = True
                    cell.auto_size = None

                    if isinstance(val,str):
                        cell.fill.fore_color.rgb = LIGHT
                    elif val <= -11:
                        cell.fill.fore_color.rgb = RGBColor(230,85,13)
                    elif val <= -4:
                        cell.fill.fore_color.rgb = RGBColor(253,141,60)
                    elif val <= 3:
                        cell.fill.fore_color.rgb = RGBColor(200,200,200)
                    elif val <= 10:
                        cell.fill.fore_color.rgb = RGBColor(66,133,244)
                    else:
                        cell.fill.fore_color.rgb = RGBColor(0,51,153)

                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        for r in p.runs:
                            r.font.size = Pt(11)

                    font_color = RGBColor(255,255,255)

                    if isinstance(val,str) or (-3 <= val <= 3):
                        font_color = RGBColor(0,0,0)

                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.color.rgb = font_color

# else:
#         # -----------------------
#     # SLIDE
#     # -----------------------


#     for col_start in range(0, len(columns), MAX_COLS):

#         col_subset = columns[col_start:col_start + MAX_COLS]

#         for row_start in range(0, len(themes), MAX_ROWS):

#             theme_subset = themes[row_start:row_start + MAX_ROWS]

#             data_subset = [
#                 row[col_start:col_start + MAX_COLS]
#                 for row in data[row_start:row_start + MAX_ROWS]
#             ]
#             slide = setup_slide(prs)
#             add_legend(slide)



#             # -----------------------
#             # TABLE
#             # -----------------------
#             rows = len(theme_subset) + 2
#             cols = len(col_subset) + 1

#             table = slide.shapes.add_table(
#                 rows, cols,
#                 Inches(0), Inches(1.5),
#                 prs.slide_width, Inches(0.3 * rows)
#             ).table

#             # COLUMN WIDTH
#             table.columns[0].width = Inches(2.5)
#             remaining = prs.slide_width - Inches(2.5)

#             for col_i in range(1, cols):
#                 table.columns[col_i].width = int(remaining / (cols - 1))

#             # HEADERS
#             # SAFE HEADERS
          
#             col_index = 1

#             for group_name, group_cols in group_definitions.items():

#                 # 🚨 STOP if no space left
#                 if col_index >= cols:
#                     break

#                 matched_cols = [c for c in col_subset if c in group_cols]

#                 if not matched_cols:
#                     continue

#                 start = col_index
#                 end = col_index + len(matched_cols) - 1

#                 # 🚨 CLAMP TO TABLE LIMIT
#                 end = min(end, cols - 1)

#                 # 🚨 DOUBLE SAFETY
#                 if start >= cols:
#                     break

#                 # SET HEADER
#                 table.cell(0, start).text = group_name

#                 # 🚨 SAFE MERGE (ONLY IF VALID & NOT SAME CELL)
#                 if start < end:
#                     try:
#                         table.cell(0, start).merge(table.cell(0, end))
#                     except:
#                         pass  # ignore already merged error

#                 col_index = end + 1


#             themes_cell = table.cell(0,0)
#             themes_cell.merge(table.cell(1,0))
#             themes_cell.text = "THEMES"

#             for j, col in enumerate(col_subset):
#                 table.cell(1, j+1).text = col

#             # HEADER STYLE
#             HEADER_YELLOW = RGBColor(255,200,0)

#             for i in range(2):
#                 for j in range(cols):
#                     if i==1 and j==0:
#                         continue

#                     cell = table.cell(i,j)
#                     cell.fill.solid()
#                     cell.fill.fore_color.rgb = HEADER_YELLOW
#                     cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

#                     for p in cell.text_frame.paragraphs:
#                         p.alignment = PP_ALIGN.CENTER
#                         for r in p.runs:
#                             r.font.bold = True
#                             r.font.size = Pt(14)
#                             r.font.color.rgb = RGBColor(0,0,0)

#             # BORDER
#             for row in table.rows:
#                 for cell in row.cells:
#                     set_border(cell)

#             # -----------------------
#             # FILL DATA  ✅ (FIXED POSITION)
#             # -----------------------
#             LIGHT = RGBColor(235,225,200)

#             for i, theme in enumerate(theme_subset):
#                 row_idx = i + 2

#                 # THEMES COLUMN
#                 cell = table.cell(row_idx,0)
#                 cell.text = theme

#                 percent = int(re.search(r'(\d+)%', theme).group(1))
#                 cell.fill.solid()
#                 cell.fill.fore_color.rgb = get_blue(percent)

#                 cell.text_frame.margin_left = Inches(0.05)

#                 for p in cell.text_frame.paragraphs:
#                     p.alignment = PP_ALIGN.LEFT
#                     for r in p.runs:
#                         r.font.color.rgb = RGBColor(255,255,255)
#                         r.font.size = Pt(12)

#                 # DATA CELLS
#                 for j, val in enumerate(data_subset[i]):
#                     cell = table.cell(row_idx, j+1)
#                     cell.text = str(val)

#                     cell.fill.solid()

#                     if isinstance(val,str):
#                         cell.fill.fore_color.rgb = LIGHT
#                     elif val <= -11:
#                         cell.fill.fore_color.rgb = RGBColor(230,85,13)
#                     elif val <= -4:
#                         cell.fill.fore_color.rgb = RGBColor(253,141,60)
#                     elif val <= 3:
#                         cell.fill.fore_color.rgb = RGBColor(200,200,200)
#                     elif val <= 10:
#                         cell.fill.fore_color.rgb = RGBColor(66,133,244)
#                     else:
#                         cell.fill.fore_color.rgb = RGBColor(0,51,153)

#                     for p in cell.text_frame.paragraphs:
#                         p.alignment = PP_ALIGN.CENTER
#                         for r in p.runs:
#                             r.font.size = Pt(11)

#                     font_color = RGBColor(255,255,255)

#                     if isinstance(val,str) or (-3 <= val <= 3):
#                         font_color = RGBColor(0,0,0)

#                     for p in cell.text_frame.paragraphs:
#                         for r in p.runs:
#                             r.font.color.rgb = font_color

# -----------------------
# SAVE
# -----------------------
prs.save("Matrix_Final_Exact.pptx")
print("✅ PPT created successfully!")